#!/usr/bin/env python3
"""
Export presentation.html to exports/VibeCoding_Presentation.pptx by screenshotting
each slide at 1920x1080 and embedding images into a 16:9 PowerPoint.

Run from anywhere:
  python3 scripts/export_html_to_pptx.py

Requires: pip install playwright python-pptx
           playwright install chromium
"""
from __future__ import annotations

import argparse
import http.server
import socketserver
import threading
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

try:
    from playwright.sync_api import sync_playwright
except ImportError as e:
    raise SystemExit(
        "Install Playwright: pip install playwright && playwright install chromium"
    ) from e


# Repository root (parent of scripts/)
ROOT = Path(__file__).resolve().parent.parent
HTML_FILE = "presentation.html"
EXPORT_DIR = ROOT / "exports"
OUTPUT_PPTX = EXPORT_DIR / "VibeCoding_Presentation.pptx"
SCREENSHOT_DIR = EXPORT_DIR / ".ppt_screenshots"
VIEWPORT = {"width": 1920, "height": 1080}
# 16:9 slide size (inches) — common widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _free_port() -> int:
    with socketserver.TCPServer(("127.0.0.1", 0), None) as s:
        return s.server_address[1]


def _make_handler_class(directory: Path):
    base = str(directory.resolve())

    class Handler(http.server.SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=base, **kwargs)

        def log_message(self, format, *args):
            pass

    return Handler


def _run_http_server(port: int, directory: Path) -> socketserver.TCPServer:
    handler = _make_handler_class(directory)
    httpd = socketserver.TCPServer(("127.0.0.1", port), handler)
    thread = threading.Thread(target=httpd.serve_forever, daemon=True)
    thread.start()
    return httpd


# Chromium args that reduce "all black" screenshots in headless (GPU/compositing issues)
_CHROMIUM_ARGS = [
    "--disable-dev-shm-usage",
    "--force-color-profile=srgb",
    "--disable-features=IsolateOrigins,site-per-process",
]


def _launch_chromium(p, prefer_system_chrome: bool):
    """Prefer system Chrome when available — often renders CSS/gradients correctly vs bundled headless."""
    if prefer_system_chrome:
        try:
            return p.chromium.launch(
                headless=True,
                channel="chrome",
                args=_CHROMIUM_ARGS,
            )
        except Exception:
            pass
    return p.chromium.launch(headless=True, args=_CHROMIUM_ARGS)


def capture_slides(
    url: str, out_dir: Path, *, prefer_system_chrome: bool = True
) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []

    with sync_playwright() as p:
        browser = _launch_chromium(p, prefer_system_chrome)
        context = browser.new_context(
            viewport=VIEWPORT,
            device_scale_factor=1,
            color_scheme="dark",
            reduced_motion="reduce",
        )
        page = context.new_page()

        page.goto(url, wait_until="domcontentloaded", timeout=120_000)
        page.wait_for_load_state("networkidle", timeout=120_000)

        # Force backgrounds to paint in headless Chromium (avoids all-black captures)
        page.add_style_tag(
            content="""
              html, body {
                background-color: #1a1a2e !important;
                -webkit-print-color-adjust: exact !important;
                print-color-adjust: exact !important;
              }
              .presentation {
                background-color: #1a1a2e !important;
                -webkit-print-color-adjust: exact !important;
                print-color-adjust: exact !important;
              }
              .slide {
                -webkit-print-color-adjust: exact !important;
                print-color-adjust: exact !important;
              }
              .slide.active {
                background-color: #16213e !important;
                background-image: linear-gradient(135deg, #1a1a2e 0%, #16213e 50%, #0f3460 100%) !important;
              }
              .slide.title-slide.active {
                background-color: #0f3460 !important;
                background-image: linear-gradient(135deg, #1a1a2e 0%, #0f3460 100%) !important;
              }
            """
        )

        page.wait_for_timeout(2000)

        total = page.evaluate("document.querySelectorAll('.slide').length")
        if total < 1:
            browser.close()
            raise RuntimeError("No .slide elements found")

        for i in range(total):
            page.evaluate(
                """(idx) => {
                    if (typeof showSlide === 'function') {
                        showSlide(idx);
                    } else {
                        const slides = document.querySelectorAll('.slide');
                        slides.forEach((s, j) => {
                            s.classList.toggle('active', j === idx);
                        });
                    }
                }""",
                i,
            )
            page.wait_for_selector(".slide.active", state="visible", timeout=10_000)
            page.wait_for_timeout(400)
            # Wait two animation frames so gradients/paints flush
            page.evaluate(
                """() => new Promise((resolve) => {
                    requestAnimationFrame(() => requestAnimationFrame(resolve));
                })"""
            )

            path = out_dir / f"slide_{i + 1:02d}.png"
            # Capture the visible slide layer (includes its background), not a blank compositor buffer
            active = page.locator(".slide.active").first
            active.screenshot(path=str(path), type="png", animations="disabled")
            paths.append(path)

        context.close()
        browser.close()

    return paths


def build_pptx(image_paths: list[Path], output: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for img in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img), 0, 0, width=SLIDE_W, height=SLIDE_H)

    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser(description="HTML presentation → PPTX via screenshots")
    parser.add_argument(
        "--keep-screenshots",
        action="store_true",
        help="Keep exports/.ppt_screenshots after building (default: delete temp PNGs)",
    )
    parser.add_argument(
        "--bundled-chromium-only",
        action="store_true",
        help="Do not use system Google Chrome (channel=chrome); fixes vary by machine",
    )
    args = parser.parse_args()

    port = _free_port()
    url = f"http://127.0.0.1:{port}/{HTML_FILE}"

    httpd = _run_http_server(port, ROOT)
    time.sleep(0.3)

    try:
        print(f"Serving {ROOT} at {url}")
        print("Capturing slides with Playwright…")
        shots = capture_slides(
            url,
            SCREENSHOT_DIR,
            prefer_system_chrome=not args.bundled_chromium_only,
        )
        print(f"Captured {len(shots)} slides")

        build_pptx(shots, OUTPUT_PPTX)
        print(f"Saved: {OUTPUT_PPTX}")
    finally:
        httpd.shutdown()

    if not args.keep_screenshots:
        for pth in SCREENSHOT_DIR.glob("slide_*.png"):
            pth.unlink(missing_ok=True)
        try:
            SCREENSHOT_DIR.rmdir()
        except OSError:
            pass


if __name__ == "__main__":
    main()
